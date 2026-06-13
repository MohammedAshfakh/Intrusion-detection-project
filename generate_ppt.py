from pptx import Presentation

ppt = Presentation()

def slide(title, content):
    layout = ppt.slide_layouts[1]
    s = ppt.slides.add_slide(layout)
    s.shapes.title.text = title
    s.placeholders[1].text = content

# ================= 10 SLIDES =================

slide("AI SOC Project", "Intrusion Detection System using ML + Flask")

slide("Problem Statement", "Cyber attacks are increasing, need real-time detection system")

slide("Objective", "Detect malicious URLs and simulate SOC monitoring")

slide("System Architecture", "Frontend + Backend + ML Model + JSON DB")

slide("Modules", "Login, Register, Dashboard, Scan, Analytics")

slide("Machine Learning", "URL feature extraction + classification")

slide("Live Threat System", "Simulated real-time attack detection")

slide("Analytics", "Fake global traffic visualization")

slide("Deployment", "Hosted on Render cloud platform")

slide("Conclusion", "SOC simulation helps understand cybersecurity systems")

ppt.save("SOC_Presentation.pptx")

print("PPT GENERATED SUCCESSFULLY")
